#!/usr/bin/env python3
"""Generate FR + AR Ministry funding PPTX dossiers for Éditions Al-Amân."""
from __future__ import annotations

import os
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# --- paths ---
MEDIA = Path("/tmp/pptx_media")
OUT_FR = Path("/home/luceor/Documents/Al-Amane_DossierFinancement_FR.pptx")
OUT_AR = Path("/home/luceor/Documents/Al-Amane_DossierFinancement_AR.pptx")
LOG = Path(__file__).with_name("progress.log")
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# --- palette (match source deck) ---
BLUE = RGBColor(0x1E, 0x4F, 0xD6)
BLUE2 = RGBColor(0x4C, 0x7B, 0xEF)
PALE = RGBColor(0xA6, 0xC5, 0xFC)
NAVY = RGBColor(0x0B, 0x1F, 0x4D)
DARK = RGBColor(0x16, 0x21, 0x3D)
GRAY = RGBColor(0x6B, 0x7A, 0x99)
MUTED = RGBColor(0x8C, 0xA1, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NOTE_BG = RGBColor(0xEA, 0xF1, 0xFE)
DEEP = RGBColor(0x15, 0x31, 0x7A)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFF)

BUDGET_ROWS = [
    (1, 3000),
    (2, 3500),
    (3, 20000),
    (4, 10000),
    (5, 10000),
    (6, 6000),
    (7, 12000),
    (8, 3500),
    (9, 700),
    (10, 12000),
    (11, 15000),
    (12, 5000),
    (13, 8000),
    (14, 4000),
    (15, 8000),
    (16, 2500),
    (17, 3800),
]
assert sum(c for _, c in BUDGET_ROWS) == 127000


def log(msg: str) -> None:
    line = f"{datetime.now().isoformat(timespec='seconds')} {msg}"
    print(line, flush=True)
    try:
        LOG.parent.mkdir(parents=True, exist_ok=True)
        with open(LOG, "a", encoding="utf-8") as f:
            f.write(line + "\n")
        with open("/tmp/funding_pptx/progress.log", "a", encoding="utf-8") as f:
            f.write(line + "\n")
    except OSError:
        pass


def fmt_tnd(n: int) -> str:
    return f"{n:,}".replace(",", " ")


def set_run(run, text, size_pt, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=DARK, font="Calibri", align=PP_ALIGN.LEFT, rtl=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT if rtl else align
    # clear any default empty run
    for r in list(p.runs):
        r.text = ""
    run = p.add_run()
    set_run(run, text, size, bold, color, font)
    if rtl:
        try:
            pPr = p._p.get_or_add_pPr()
            pPr.set(qn("a:rtl"), "1")
        except Exception:
            pass
    return box


def text_block(slide, left, top, width, height, lines, rtl=False):
    """lines: list of (text, size, bold, color, font?)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, size, bold, color = item[0], item[1], item[2], item[3]
        font = item[4] if len(item) > 4 else "Calibri"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
        run = p.add_run()
        set_run(run, text, size, bold, color, font)
        if rtl:
            try:
                p._p.get_or_add_pPr().set(qn("a:rtl"), "1")
            except Exception:
                pass
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def footer(slide, num, brand, rtl=False):
    logo = MEDIA / "logo.jpg"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.5), Inches(6.95), Inches(0.34), Inches(0.34))
    add_textbox(slide, Inches(0.92), Inches(6.98), Inches(3), Inches(0.3),
                brand, size=9, color=GRAY, rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)
    add_textbox(slide, Inches(12.4), Inches(6.98), Inches(0.6), Inches(0.3),
                f"{num:02d}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def decor_corner(slide):
    add_oval(slide, Inches(11.35), Inches(-2.3), Inches(4.6), Inches(4.6), PALE)
    add_oval(slide, Inches(12.75), Inches(1.6), Inches(1.6), Inches(1.6), BLUE)


def header(slide, eyebrow, title, rtl=False):
    decor_corner(slide)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(10), Inches(0.35),
                eyebrow, size=12, bold=True, color=BLUE, rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.65), Inches(0.75), Inches(11), Inches(0.7),
                title, size=28, bold=True, color=NAVY, font="Cambria", rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)


def card(slide, left, top, width, height, title, body, accent=BLUE, rtl=False):
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.08), height, accent)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.4), Inches(0.4),
                title, size=13, bold=True, color=NAVY, rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.6), width - Inches(0.4), height - Inches(0.75),
                body, size=11, color=DARK, rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)


def note_bar(slide, text, rtl=False):
    add_rect(slide, Inches(0.65), Inches(6.0), Inches(11.9), Inches(0.7), NOTE_BG)
    add_textbox(slide, Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.5),
                text, size=11, color=DEEP, rtl=rtl,
                align=PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)


def table_slide(slide, rows, left, top, width, col_widths, rtl=False):
    """rows: list of list of str. First row = header."""
    n_rows, n_cols = len(rows), len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.38 * n_rows))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if (rtl or c > 0) else PP_ALIGN.LEFT
            if c == 0 and not rtl:
                p.alignment = PP_ALIGN.LEFT
            if rtl:
                p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            is_header = r == 0
            is_total = r == n_rows - 1 and any("127" in str(x) or "الإجمالي" in str(x) or "TOTAL" in str(x) or "إجمالي" in str(x) for x in row)
            set_run(run, str(val), 10 if n_rows > 10 else 11, bold=is_header or is_total,
                    color=WHITE if is_header else NAVY)
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# ---------- content strings ----------
C = {
    "fr": {
        "brand": "ÉDITIONS AL-AMÂN",
        "brand_ar_line": "دار أمان الله للنشر والتوزيع",
        "rtl": False,
        "cover_sub": "Dossier de candidature — Appel à projets",
        "cover_title": "Al-Amân",
        "cover_tag": "Création d'une plateforme éducative culturelle numérique\nbilingue arabe / français",
        "cover_budget": "Budget demandé : 100 000 TND  ·  Total projet : 127 000 TND",
        "cover_ministry": "Ministère des Affaires Culturelles — Portail des subventions",
        "cover_status": "Statut : Phase 1 en cours de développement",
        "toc_eye": "SOMMAIRE",
        "toc_title": "Table des matières",
        "toc": [
            "1. Page de couverture",
            "2. Table des matières",
            "3. Fiche d'identité du projet",
            "4. Résumé exécutif",
            "5. Problématique et contexte",
            "6. Opportunité",
            "7. Solution et profils utilisateurs",
            "8. Fonctionnalités principales",
            "9. État d'avancement (en développement)",
            "10. Étude de faisabilité",
            "11. Étude technique",
            "12. Sécurité et fiabilité",
            "13. Plan de contenu et numérisation",
            "14–15. Budget détaillé (127 000 TND)",
            "16. Plan de financement",
            "17. Contribution de la porteuse",
            "18. Ressources humaines",
            "19. Calendrier (12 mois)",
            "20. Analyse des risques",
            "21. Plan de durabilité",
            "22. Indicateurs de succès",
            "23–27. Annexes (captures & logo)",
            "28. Conclusion",
        ],
        "id_eye": "FICHE PROJET",
        "id_title": "Carte d'identité du projet",
        "id_rows": [
            ("Titre", "Création d'une plateforme éducative culturelle numérique"),
            ("Porteur", "Éditions Al-Amân — دار أمان الله للنشر والتوزيع"),
            ("Durée", "12 mois"),
            ("Coût total", "127 000 TND"),
            ("Soutien demandé", "100 000 TND (79 %) — Ministère des Affaires Culturelles"),
            ("Contribution propre", "27 000 TND (21 %)"),
            ("Statut", "Phase 1 — en cours de développement"),
            ("Périmètre", "Web, Android, iOS, CMS, IA, numérisation, production média"),
            ("Contact", "contact@al-amane-edu.tn"),
        ],
        "sum_eye": "VUE D'ENSEMBLE",
        "sum_title": "Résumé exécutif",
        "sum_cards": [
            ("Une plateforme bilingue", "Al-Amân s'adresse aux élèves du primaire, enseignants, parents et administrateurs, avec des contenus éducatifs et culturels en arabe et en français."),
            ("Un besoin documenté", "Les ressources numériques tunisiennes sont fragmentées (YouTube, Facebook, WhatsApp) sans structure pédagogique ni éditeur culturel unifié."),
            ("Un projet déjà en marche", "Phase 1 en cours : la plateforme est fonctionnelle en local (auth, catalogue, admin, bilingue). Le développement se poursuit."),
            ("Un financement ciblé", "100 000 TND de soutien ministériel + 27 000 TND de contribution propre pour finaliser Web/mobile/IA, contenu et lancement."),
        ],
        "sum_note": "Statut : Phase 1 en développement — le financement finalise la plateforme, le contenu culturel et le lancement public.",
        "prob_eye": "LE CONTEXTE",
        "prob_title": "Problématique",
        "prob_intro": "En Tunisie, l'éducation et la culture numériques au primaire manquent de structuration, malgré une forte demande des familles et des enseignants.",
        "prob_items": [
            ("Ressources dispersées", "Contenus éparpillés sans cohérence pédagogique ni éditoriale."),
            ("Aucune plateforme centralisée", "Pas de solution unique dédiée au primaire tunisien bilingue."),
            ("Contenu bilingue insuffisant", "Peu de ressources combinent arabe et français de façon cohérente."),
            ("Suivi parental limité", "Peu d'outils pour suivre la progression des enfants."),
            ("Manque d'outils enseignants", "Peu de supports interactifs prêts pour la classe."),
            ("Fracture numérique", "Accès inégal selon les régions et les foyers."),
        ],
        "opp_eye": "POURQUOI MAINTENANT",
        "opp_title": "Opportunité",
        "opp_stats": [
            ("1M+", "élèves scolarisés dans le primaire en Tunisie"),
            ("80%+", "de pénétration smartphone dans les foyers"),
            ("Maison d'édition", "fonds culturel et éducatif déjà existant à numériser"),
            ("0", "concurrent direct bilingue structuré pour l'école publique"),
        ],
        "opp_body": "Le Ministère encourage la transformation numérique de la culture et de l'éducation. Al-Amân se situe à l'intersection d'une maison d'édition tunisienne, d'une plateforme bilingue et d'un accès national via le web et le mobile.",
        "sol_eye": "LE PRODUIT",
        "sol_title": "Notre solution — quatre profils",
        "sol_roles": [
            ("É", "Élève", "Cours vidéo, livres PDF, exercices adaptés, suivi de progression."),
            ("E", "Enseignant", "Ressources pédagogiques et outils numériques pour la classe."),
            ("P", "Parent", "Suivi de la progression et gestion de l'abonnement familial."),
            ("A", "Admin", "Gestion des contenus, utilisateurs et paiements."),
        ],
        "feat_eye": "LE PRODUIT",
        "feat_title": "Fonctionnalités principales",
        "feat_cols": [
            ("Accès & comptes", ["Authentification multi-rôles", "OAuth Google", "Récupération de mot de passe"]),
            ("Contenu", ["Vidéos, PDF, animations", "Catalogue & recherche", "Interface AR/FR (RTL/LTR)"]),
            ("Plateformes", ["Application Web", "Android (en cours)", "iOS (prévu)"]),
            ("Évolutions financées", ["Assistant IA pédagogique", "CMS & base de données", "Paiements & lancement"]),
        ],
        "dev_eye": "PREUVE D'EXÉCUTION",
        "dev_title": "État d'avancement — toujours en développement",
        "dev_done": [
            "Authentification à 4 rôles opérationnelle",
            "Catalogue de cours avec filtres et recherche",
            "Lecture vidéo (gratuit / premium)",
            "Interface bilingue arabe / français",
            "Sécurité de base (rate limit, en-têtes)",
            "Tableau de bord administrateur",
            "Application Android Expo (base)",
            "Assistant IA (prototype)",
        ],
        "dev_todo": [
            "Finaliser Web + CMS + base de production",
            "Applications Android & iOS complètes",
            "Intégration IA pédagogique avancée",
            "Numérisation de 500+ ouvrages",
            "Production de 100+ leçons / vidéos",
            "Hébergement cloud & domaine .tn",
            "Marketing et lancement public",
            "Partenariats écoles & institutions",
        ],
        "dev_note": "Le projet est fonctionnel en local mais reste en phase de développement. Ce dossier finance la finalisation et le déploiement public.",
        "fea_eye": "ÉTUDE DE FAISABILITÉ",
        "fea_title": "Faisabilité du projet",
        "fea_cards": [
            ("Demande réelle", "Familles et enseignants cherchent un contenu structuré, bilingue et fiable."),
            ("Actif éditorial", "دار أمان الله dispose déjà de droits et de fonds à numériser."),
            ("Base technique", "Architecture moderne déjà amorcée (Next.js, mobile, IA)."),
            ("Modèle durable", "Abonnements, cours, livres numériques et partenariats scolaires."),
        ],
        "tech_eye": "ÉTUDE TECHNIQUE",
        "tech_title": "Architecture technique",
        "tech_items": [
            ("1", "Frontend Web", "Next.js / TypeScript, responsive, bilingue AR/FR."),
            ("2", "Applications mobiles", "Android (Expo) en cours ; iOS dans le plan de 12 mois."),
            ("3", "Backend & CMS", "API sécurisée, gestion de contenu et rôles."),
            ("4", "Données & cloud", "PostgreSQL, stockage fichiers, streaming vidéo, CDN."),
            ("5", "Intelligence artificielle", "Assistant pédagogique (recherche, explications, aide)."),
            ("6", "Exploitation", "Emails, monitoring, CI/CD, certificats SSL, domaine."),
        ],
        "sec_eye": "CONFIANCE",
        "sec_title": "Sécurité et fiabilité",
        "sec_items": [
            "Contrôle d'accès par rôle (élève, enseignant, parent, admin)",
            "Chiffrement des mots de passe",
            "Limitation du taux de requêtes",
            "En-têtes de sécurité web",
            "Journalisation et audit des actions sensibles",
            "Validation stricte des données côté serveur",
        ],
        "cont_eye": "CONTENU CULTUREL",
        "cont_title": "Plan de contenu et numérisation",
        "cont_steps": [
            ("1", "Script pédagogique"),
            ("2", "Validation éditoriale"),
            ("3", "Numérisation / tournage"),
            ("4", "Montage & voix off"),
            ("5", "Relecture bilingue"),
            ("6", "Publication"),
        ],
        "cont_stats": [
            ("500+", "livres et publications numériques"),
            ("100+", "leçons et vidéos éducatives"),
            ("AR/FR", "publication bilingue"),
            ("12 mois", "cycle de production"),
        ],
        "bud_eye": "ÉTUDE FINANCIÈRE",
        "bud1_title": "Budget détaillé (1/2)",
        "bud2_title": "Budget détaillé (2/2) — total 127 000 TND",
        "bud_headers": ["N°", "Poste", "Coût (TND)"],
        "bud_labels": [
            "Études techniques et analyse des besoins",
            "Identité visuelle, logo et interfaces",
            "Développement de la plateforme Web",
            "Développement application Android",
            "Développement application iOS",
            "CMS et base de données",
            "Intégration IA et assistant intelligent",
            "Hébergement cloud (1 an)",
            "Nom de domaine et certificats SSL",
            "Numérisation des livres et publications",
            "Production vidéos, leçons et médias",
            "Voix off et montage",
            "Marketing, publicité et lancement",
            "Formation de l'équipe",
            "Gestion, suivi et évaluation",
            "Frais juridiques et administratifs",
            "Réserve pour risques et imprévus",
        ],
        "fund_eye": "PLAN DE FINANCEMENT",
        "fund_title": "Sources de financement",
        "fund_headers": ["Source", "Montant (TND)", "Part"],
        "fund_rows": [
            ["Contribution de la porteuse du projet", "27 000", "21 %"],
            ["Soutien demandé — Ministère des Affaires Culturelles", "100 000", "79 %"],
            ["TOTAL", "127 000", "100 %"],
        ],
        "own_eye": "CONTRIBUTION PROPRE",
        "own_title": "Contribution de la porteuse du projet",
        "own_headers": ["Contribution", "Valeur (TND)"],
        "own_rows": [
            ["Équipements informatiques", "8 000"],
            ["Droits d'édition et contenu scientifique", "9 000"],
            ["Équipements de bureau", "3 000"],
            ["Travail administratif et direction de projet", "5 000"],
            ["Apport financier direct", "2 000"],
            ["TOTAL", "27 000"],
        ],
        "hr_eye": "RESSOURCES HUMAINES",
        "hr_title": "Équipe projet",
        "hr_headers": ["Spécialité", "Effectif"],
        "hr_rows": [
            ["Directeur / Directrice de projet", "1"],
            ["Directeur technique", "1"],
            ["Développeur Full Stack", "1"],
            ["Développeur applications mobiles", "1"],
            ["Designer UX/UI", "1"],
            ["Spécialiste contenu pédagogique", "2"],
            ["Spécialiste production audiovisuelle", "1"],
            ["Spécialiste marketing digital", "1"],
            ["Support technique & service clients", "1"],
        ],
        "tl_eye": "MISE EN ŒUVRE",
        "tl_title": "Calendrier — 12 mois",
        "tl_phases": [
            ("Mois 1–2", "Études & design", "Analyse des besoins, identité visuelle, maquettes UX/UI."),
            ("Mois 3–5", "Web & CMS", "Finalisation plateforme Web, CMS, base de données."),
            ("Mois 5–8", "Mobile & IA", "Android, iOS, intégration de l'assistant intelligent."),
            ("Mois 6–10", "Contenu", "Numérisation, vidéos, voix off, validation pédagogique."),
            ("Mois 10–12", "Lancement", "Hébergement, formation, marketing, évaluation."),
        ],
        "risk_eye": "RISQUES",
        "risk_title": "Analyse des risques",
        "risk_items": [
            ("Technique", "Retards de développement", "Équipe déjà en place ; jalons mensuels ; réserve 3 800 TND."),
            ("Contenu", "Retard de production éditoriale", "Pipeline validé avec enseignants ; lots prioritaires."),
            ("Adoption", "Adoption plus lente", "Partenariats écoles ; offre freemium ; campagne ciblée."),
            ("Financier", "Décalage du soutien", "Contribution propre 27 000 TND déjà mobilisable."),
        ],
        "sus_eye": "DURABILITÉ",
        "sus_title": "Plan de durabilité",
        "sus_cards": [
            ("Abonnements", "Revenus récurrents élèves / parents / enseignants (TND)."),
            ("Livres numériques", "Monétisation du catalogue numérisé de la maison d'édition."),
            ("Cours & formations", "Modules payants et packs institutionnels."),
            ("Partenariats", "Écoles, institutions culturelles et collectives."),
        ],
        "kpi_eye": "INDICATEURS",
        "kpi_title": "Indicateurs de succès",
        "kpi_items": [
            ("500+", "livres et publications culturelles numérisés"),
            ("100+", "leçons et vidéos éducatives produites"),
            ("5 000", "utilisateurs la première année"),
            ("Partenariats", "écoles et institutions culturelles"),
            ("5", "emplois directs créés (+ indirects)"),
            ("Revenus", "abonnements, cours, livres numériques"),
        ],
        "ax_eye": "ANNEXES",
        "ax_home": "Captures — page d'accueil",
        "ax_admin": "Captures — administration",
        "ax_student": "Captures — espace élève",
        "ax_teacher": "Captures — espace enseignant",
        "ax_logo": "Annexe — logo دار أمان الله للنشر والتوزيع",
        "close_title": "Merci de l'attention portée à ce dossier",
        "close_budget": "Budget demandé : 100 000 TND",
        "close_total": "Coût total du projet : 127 000 TND",
        "close_status": "Statut : Phase 1 en cours de développement · Plateforme déjà fonctionnelle en local",
        "close_portal": "portail.mac-subventions.gov.tn",
        "close_contact": "Contact : contact@al-amane-edu.tn",
        "close_min": "Ministère des Affaires Culturelles — Dossier de candidature",
    },
    "ar": {
        "brand": "دار أمان الله",
        "brand_ar_line": "دار أمان الله للنشر والتوزيع",
        "rtl": True,
        "cover_sub": "ملف ترشح — دعوة لتقديم المشاريع",
        "cover_title": "أمان الله",
        "cover_tag": "إنشاء منصة تعليمية ثقافية رقمية\nثنائية اللغة عربي / فرنسي",
        "cover_budget": "الدعم المطلوب: 100 000 د.ت  ·  إجمالي المشروع: 127 000 د.ت",
        "cover_ministry": "وزارة الشؤون الثقافية — بوابة المنح",
        "cover_status": "الحالة: المرحلة 1 قيد التطوير",
        "toc_eye": "الفهرس",
        "toc_title": "فهرس المحتويات",
        "toc": [
            "1. صفحة الغلاف الرسمية",
            "2. فهرس المحتويات",
            "3. بطاقة تعريف بالمشروع",
            "4. الملخص التنفيذي",
            "5. الإشكالية والسياق",
            "6. الفرصة",
            "7. الحل وملفات المستخدمين",
            "8. الوظائف الرئيسية",
            "9. حالة التقدم (قيد التطوير)",
            "10. دراسة الجدوى",
            "11. الدراسة التقنية",
            "12. الأمن والموثوقية",
            "13. خطة المحتوى والرقمنة",
            "14–15. الميزانية التفصيلية (127 000 د.ت)",
            "16. خطة التمويل",
            "17. مساهمة صاحبة المشروع",
            "18. الموارد البشرية",
            "19. الرزنامة (12 شهراً)",
            "20. تحليل المخاطر",
            "21. خطة الاستدامة",
            "22. مؤشرات النجاح",
            "23–27. الملاحق (لقطات وشعار)",
            "28. الخاتمة",
        ],
        "id_eye": "بطاقة المشروع",
        "id_title": "بطاقة تعريف بالمشروع",
        "id_rows": [
            ("العنوان", "إنشاء منصة تعليمية ثقافية رقمية"),
            ("الجهة", "دار أمان الله للنشر والتوزيع"),
            ("مدة الإنجاز", "12 شهراً"),
            ("التكلفة الإجمالية", "127 000 د.ت"),
            ("الدعم المطلوب", "100 000 د.ت (79 %) — وزارة الشؤون الثقافية"),
            ("المساهمة الذاتية", "27 000 د.ت (21 %)"),
            ("الحالة", "المرحلة 1 — قيد التطوير"),
            ("النطاق", "ويب، أندرويد، iOS، إدارة محتوى، ذكاء اصطناعي، رقمنة، إنتاج وسائط"),
            ("التواصل", "contact@al-amane-edu.tn"),
        ],
        "sum_eye": "نظرة عامة",
        "sum_title": "الملخص التنفيذي",
        "sum_cards": [
            ("منصة ثنائية اللغة", "تستهدف أمان الله تلاميذ الابتدائي والمعلمين والأولياء والإدارة، بمحتوى تعليمي وثقافي بالعربية والفرنسية."),
            ("حاجة موثقة", "الموارد الرقمية التونسية مشتتة (يوتيوب، فيسبوك، واتساب) دون هيكلة تربوية أو دار نشر رقمية موحدة."),
            ("مشروع قيد الإنجاز", "المرحلة 1 جارية: المنصة تعمل محلياً (مصادقة، كتالوج، إدارة، ثنائية اللغة). التطوير مستمر."),
            ("تمويل موجّه", "100 000 د.ت دعماً وزارياً + 27 000 د.ت مساهمة ذاتية لإتمام الويب والجوال والذكاء الاصطناعي والمحتوى والإطلاق."),
        ],
        "sum_note": "الحالة: المرحلة 1 قيد التطوير — التمويل يُتمّ المنصة والمحتوى الثقافي والإطلاق العمومي.",
        "prob_eye": "السياق",
        "prob_title": "الإشكالية",
        "prob_intro": "في تونس، يعاني التعليم والثقافة الرقمية في الابتدائي من ضعف الهيكلة رغم الطلب القوي من العائلات والمعلمين.",
        "prob_items": [
            ("موارد مشتتة", "محتويات مبعثرة دون اتساق تربوي أو تحريري."),
            ("لا منصة مركزية", "لا حل موحّد مخصص للابتدائي التونسي ثنائي اللغة."),
            ("محتوى ثنائي محدود", "قليل من الموارد يجمع العربية والفرنسية بشكل متماسك."),
            ("متابعة أسرية محدودة", "أدوات قليلة لمتابعة تقدّم الأبناء."),
            ("نقص أدوات المعلمين", "قليل من الدعائم التفاعلية الجاهزة للقسم."),
            ("فجوة رقمية", "نفاذ غير متكافئ حسب الجهات والأسر."),
        ],
        "opp_eye": "لماذا الآن",
        "opp_title": "الفرصة",
        "opp_stats": [
            ("+1M", "تلميذ في الابتدائي بتونس"),
            ("+%80", "انتشار الهواتف الذكية في الأسر"),
            ("دار نشر", "رصيد ثقافي وتعليمي جاهز للرقمنة"),
            ("0", "منافس مباشر ثنائي اللغة منظم للمدرسة العمومية"),
        ],
        "opp_body": "تشجّع الوزارة التحوّل الرقمي للثقافة والتعليم. تقع أمان الله عند تقاطع دار نشر تونسية ومنصة ثنائية اللغة ونفاذ وطني عبر الويب والجوال.",
        "sol_eye": "المنتج",
        "sol_title": "الحل — أربعة ملفات",
        "sol_roles": [
            ("ت", "التلميذ", "دروس فيديو وكتب PDF وتمارين ومتابعة التقدّم."),
            ("م", "المعلم", "موارد تربوية وأدوات رقمية للقسم."),
            ("و", "الولي", "متابعة التقدّم وإدارة الاشتراك العائلي."),
            ("إ", "الإدارة", "إدارة المحتويات والمستخدمين والمدفوعات."),
        ],
        "feat_eye": "المنتج",
        "feat_title": "الوظائف الرئيسية",
        "feat_cols": [
            ("الوصول والحسابات", ["مصادقة متعددة الأدوار", "Google OAuth", "استعادة كلمة المرور"]),
            ("المحتوى", ["فيديو وPDF ورسوم", "كتالوج وبحث", "واجهة عربي/فرنسي"]),
            ("المنصات", ["تطبيق ويب", "أندرويد (جارٍ)", "iOS (مخطط)"]),
            ("ما يموّله الدعم", ["مساعد ذكي تربوي", "نظام إدارة محتوى", "مدفوعات وإطلاق"]),
        ],
        "dev_eye": "إثبات الإنجاز",
        "dev_title": "حالة التقدم — ما زال قيد التطوير",
        "dev_done": [
            "مصادقة بأربعة أدوار جاهزة",
            "كتالوج دروس مع تصفية وبحث",
            "تشغيل فيديو (مجاني / مدفوع)",
            "واجهة ثنائية عربي / فرنسي",
            "أمن أساسي (حد المعدّل والعناوين)",
            "لوحة تحكم إدارية",
            "تطبيق أندرويد Expo (أساس)",
            "مساعد ذكاء اصطناعي (نموذج أولي)",
        ],
        "dev_todo": [
            "إتمام الويب ونظام المحتوى وقاعدة الإنتاج",
            "تطبيقات أندرويد وiOS كاملة",
            "دمج الذكاء الاصطناعي التربوي المتقدم",
            "رقمنة أكثر من 500 مؤلف",
            "إنتاج أكثر من 100 درس / فيديو",
            "استضافة سحابية ونطاق .tn",
            "تسويق وإطلاق عمومي",
            "شراكات مع مدارس ومؤسسات",
        ],
        "dev_note": "المشروع يعمل محلياً لكنه ما زال في طور التطوير. يموّل هذا الملف الإتمام والنشر العمومي.",
        "fea_eye": "دراسة الجدوى",
        "fea_title": "جدوى المشروع",
        "fea_cards": [
            ("طلب حقيقي", "العائلات والمعلمون يبحثون عن محتوى منظم وثنائي وموثوق."),
            ("رصيد تحريري", "لدار أمان الله حقوق ومحتوى جاهز للرقمنة."),
            ("أساس تقني", "بنية حديثة بدأت فعلاً (Next.js، جوّال، ذكاء اصطناعي)."),
            ("نموذج مستدام", "اشتراكات ودروس وكتب رقمية وشراكات مدرسية."),
        ],
        "tech_eye": "الدراسة التقنية",
        "tech_title": "الهندسة التقنية",
        "tech_items": [
            ("1", "واجهة الويب", "Next.js / TypeScript، متجاوبة، ثنائية اللغة."),
            ("2", "تطبيقات الجوّال", "أندرويد (Expo) جارٍ؛ iOS ضمن خطة 12 شهراً."),
            ("3", "الخادم وإدارة المحتوى", "واجهة برمجية آمنة وإدارة أدوار."),
            ("4", "البيانات والسحابة", "PostgreSQL وتخزين وبث فيديو وCDN."),
            ("5", "الذكاء الاصطناعي", "مساعد تربوي (بحث وشرح ومساعدة)."),
            ("6", "التشغيل", "بريد ومراقبة وCI/CD وشهادات ونطاق."),
        ],
        "sec_eye": "الثقة",
        "sec_title": "الأمن والموثوقية",
        "sec_items": [
            "تحكم بالوصول حسب الدور",
            "تشفير كلمات المرور",
            "الحد من معدّل الطلبات",
            "عناوين أمن الويب",
            "تسجيل ومراجعة الإجراءات الحساسة",
            "تحقق صارم من البيانات على الخادم",
        ],
        "cont_eye": "المحتوى الثقافي",
        "cont_title": "خطة المحتوى والرقمنة",
        "cont_steps": [
            ("1", "سيناريو تربوي"),
            ("2", "مصادقة تحريرية"),
            ("3", "رقمنة / تصوير"),
            ("4", "مونتاج وتعليق صوتي"),
            ("5", "مراجعة ثنائية"),
            ("6", "نشر"),
        ],
        "cont_stats": [
            ("+500", "كتاب ومنشور رقمي"),
            ("+100", "درس وفيديو تعليمي"),
            ("ع/ف", "نشر ثنائي اللغة"),
            ("12 شهراً", "دورة الإنتاج"),
        ],
        "bud_eye": "الدراسة المالية",
        "bud1_title": "الميزانية التفصيلية (1/2)",
        "bud2_title": "الميزانية التفصيلية (2/2) — الإجمالي 127 000 د.ت",
        "bud_headers": ["رقم", "البند", "الكلفة (د.ت)"],
        "bud_labels": [
            "الدراسات الفنية وتحليل الاحتياجات",
            "تصميم الهوية البصرية والشعار والواجهات",
            "تطوير المنصة الإلكترونية (Web)",
            "تطوير تطبيق Android",
            "تطوير تطبيق iOS",
            "نظام إدارة المحتوى وقاعدة البيانات",
            "دمج خدمات الذكاء الاصطناعي والمساعد الذكي",
            "استضافة سحابية وخادم لمدة سنة",
            "شراء اسم النطاق وشهادات الحماية",
            "رقمنة الكتب والمنشورات الثقافية",
            "إنتاج الفيديوهات والدروس والوسائط التعليمية",
            "تسجيل التعليق الصوتي والمونتاج",
            "التسويق والإشهار وإطلاق المنصة",
            "التدريب والتكوين لفريق العمل",
            "الإدارة والمتابعة والتقييم",
            "المصاريف القانونية والإدارية",
            "احتياطي للمخاطر والطوارئ",
        ],
        "fund_eye": "خطة التمويل",
        "fund_title": "مصادر التمويل",
        "fund_headers": ["مصدر التمويل", "القيمة (د.ت)", "النسبة"],
        "fund_rows": [
            ["مساهمة صاحبة المشروع", "27 000", "21 %"],
            ["الدعم المطلوب من وزارة الشؤون الثقافية", "100 000", "79 %"],
            ["الإجمالي", "127 000", "100 %"],
        ],
        "own_eye": "المساهمة الذاتية",
        "own_title": "مساهمة صاحبة المشروع",
        "own_headers": ["المساهمة", "القيمة (د.ت)"],
        "own_rows": [
            ["تجهيزات إعلامية", "8 000"],
            ["حقوق نشر ومحتوى علمي", "9 000"],
            ["تجهيزات مكتبية", "3 000"],
            ["عمل إداري وإدارة المشروع", "5 000"],
            ["مساهمة مالية مباشرة", "2 000"],
            ["الإجمالي", "27 000"],
        ],
        "hr_eye": "الموارد البشرية",
        "hr_title": "فريق المشروع",
        "hr_headers": ["الاختصاص", "العدد"],
        "hr_rows": [
            ["مدير المشروع", "1"],
            ["مدير تقني", "1"],
            ["مطور Full Stack", "1"],
            ["مطور تطبيقات", "1"],
            ["مصمم UX/UI", "1"],
            ["مختص محتوى تربوي", "2"],
            ["مختص إنتاج سمعي بصري", "1"],
            ["مختص تسويق رقمي", "1"],
            ["دعم فني وخدمة الحرفاء", "1"],
        ],
        "tl_eye": "التنفيذ",
        "tl_title": "الرزنامة — 12 شهراً",
        "tl_phases": [
            ("الشهر 1–2", "دراسات وتصميم", "تحليل الاحتياجات والهوية والواجهات."),
            ("الشهر 3–5", "الويب ونظام المحتوى", "إتمام المنصة ونظام الإدارة وقاعدة البيانات."),
            ("الشهر 5–8", "الجوّال والذكاء", "أندرويد وiOS ودمج المساعد الذكي."),
            ("الشهر 6–10", "المحتوى", "رقمنة وفيديوهات وتعليق صوتي ومصادقة تربوية."),
            ("الشهر 10–12", "الإطلاق", "استضافة وتدريب وتسويق وتقييم."),
        ],
        "risk_eye": "المخاطر",
        "risk_title": "تحليل المخاطر",
        "risk_items": [
            ("تقني", "تأخر التطوير", "فريق قائم؛ محطات شهرية؛ احتياطي 3 800 د.ت."),
            ("محتوى", "تأخر الإنتاج", "مسار معتمد مع معلمين؛ دفعات أولوية."),
            ("تبني", "تبني أبطأ", "شراكات مدارس؛ عرض مجاني جزئي؛ حملة موجّهة."),
            ("مالي", "تأخر الدعم", "مساهمة ذاتية 27 000 د.ت قابلة للتعبئة."),
        ],
        "sus_eye": "الاستدامة",
        "sus_title": "خطة الاستدامة",
        "sus_cards": [
            ("الاشتراكات", "إيرادات متكررة للتلاميذ والأولياء والمعلمين."),
            ("الكتب الرقمية", "تحقيق دخل من كتالوج الدار المرقمن."),
            ("الدورات", "وحدات مدفوعة وباقات مؤسساتية."),
            ("الشراكات", "مدارس ومؤسسات ثقافية وجماعات."),
        ],
        "kpi_eye": "المؤشرات",
        "kpi_title": "مؤشرات النجاح",
        "kpi_items": [
            ("+500", "كتاب ومنشور ثقافي مرقمن"),
            ("+100", "درس وفيديو تعليمي"),
            ("5 000", "مستخدم خلال السنة الأولى"),
            ("شراكات", "مع مدارس ومؤسسات ثقافية"),
            ("5", "مواطن شغل مباشرة (+ غير مباشرة)"),
            ("إيرادات", "اشتراكات ودورات وكتب رقمية"),
        ],
        "ax_eye": "الملاحق",
        "ax_home": "لقطات — الصفحة الرئيسية",
        "ax_admin": "لقطات — لوحة الإدارة",
        "ax_student": "لقطات — فضاء التلميذ",
        "ax_teacher": "لقطات — فضاء المعلم",
        "ax_logo": "ملحق — شعار دار أمان الله للنشر والتوزيع",
        "close_title": "شكراً على اهتمامكم بهذا الملف",
        "close_budget": "الدعم المطلوب: 100 000 د.ت",
        "close_total": "إجمالي تكلفة المشروع: 127 000 د.ت",
        "close_status": "الحالة: المرحلة 1 قيد التطوير · المنصة تعمل محلياً",
        "close_portal": "portail.mac-subventions.gov.tn",
        "close_contact": "التواصل: contact@al-amane-edu.tn",
        "close_min": "وزارة الشؤون الثقافية — ملف الترشح",
    },
}


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    # use blank layout if available else layout 0
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)



def build(lang: str, out: Path):
    t = C[lang]
    rtl = t["rtl"]
    log(f"BUILD start lang={lang} -> {out}")
    prs = new_prs()
    n = 0

    def next_slide():
        nonlocal n
        n += 1
        return blank(prs)

    # --- 1 Cover ---
    s = next_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_oval(s, Inches(9.6), Inches(-2.2), Inches(6.5), Inches(6.5), BLUE)
    add_oval(s, Inches(-2.2), Inches(4.8), Inches(5.5), Inches(5.5), BLUE2)
    add_oval(s, Inches(10.9), Inches(4.6), Inches(2), Inches(2), PALE)
    if (MEDIA / "logo.jpg").exists():
        s.shapes.add_picture(str(MEDIA / "logo.jpg"), Inches(10.55), Inches(0.7), Inches(1.85), Inches(1.85))
    al = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    add_textbox(s, Inches(0.8), Inches(0.7), Inches(8), Inches(0.35), t["brand"], 12, True, PALE, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(1.05), Inches(9), Inches(0.3), t["cover_sub"], 11, False, MUTED, rtl=rtl, align=al)
    add_textbox(s, Inches(0.75), Inches(2.4), Inches(10), Inches(1.0), t["cover_title"], 54, True, WHITE, "Cambria", rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(3.6), Inches(9), Inches(1.0), t["cover_tag"], 16, False, PALE, rtl=rtl, align=al)
    add_rect(s, Inches(0.8), Inches(5.0), Inches(7.2), Inches(0.5), WHITE)
    add_textbox(s, Inches(0.9), Inches(5.05), Inches(7.0), Inches(0.4), t["cover_budget"], 11, True, BLUE, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(5.65), Inches(9), Inches(0.3), t["cover_ministry"], 11, False, MUTED, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(6.0), Inches(9), Inches(0.3), t["cover_status"], 12, True, AMBER, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} cover")

    # --- 2 TOC ---
    s = next_slide()
    header(s, t["toc_eye"], t["toc_title"], rtl)
    mid = (len(t["toc"]) + 1) // 2
    for col, items in enumerate([t["toc"][:mid], t["toc"][mid:]]):
        y = Inches(1.65)
        x = Inches(0.8 + col * 6.0)
        for item in items:
            add_textbox(s, x, y, Inches(5.5), Inches(0.32), item, 12, False, DARK, rtl=rtl, align=al)
            y += Inches(0.34)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} toc")

    # --- 3 Identity ---
    s = next_slide()
    header(s, t["id_eye"], t["id_title"], rtl)
    id_header = ["Champ", "Valeur"] if not rtl else ["البند", "القيمة"]
    id_data = [id_header] + [[a, b] for a, b in t["id_rows"]]
    table_slide(s, id_data, Inches(0.7), Inches(1.6), Inches(11.8), [Inches(3.5), Inches(8.3)], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} identity")

    # --- 4 Executive summary ---
    s = next_slide()
    header(s, t["sum_eye"], t["sum_title"], rtl)
    positions = [
        (Inches(0.65), Inches(1.6)), (Inches(6.7), Inches(1.6)),
        (Inches(0.65), Inches(3.7)), (Inches(6.7), Inches(3.7)),
    ]
    accents = [BLUE, BLUE2, DEEP, AMBER]
    for (title, body), (x, y), acc in zip(t["sum_cards"], positions, accents):
        card(s, x, y, Inches(5.7), Inches(1.9), title, body, acc, rtl)
    note_bar(s, t["sum_note"], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} summary")

    # --- 5 Problem ---
    s = next_slide()
    header(s, t["prob_eye"], t["prob_title"], rtl)
    add_textbox(s, Inches(0.7), Inches(1.55), Inches(11.8), Inches(0.5), t["prob_intro"], 13, False, DARK, rtl=rtl, align=al)
    for i, (title, body) in enumerate(t["prob_items"]):
        col, row = i % 3, i // 3
        x = Inches(0.65 + col * 4.1)
        y = Inches(2.2 + row * 1.7)
        card(s, x, y, Inches(3.9), Inches(1.5), title, body, BLUE if i % 2 == 0 else BLUE2, rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} problem")

    # --- 6 Opportunity ---
    s = next_slide()
    header(s, t["opp_eye"], t["opp_title"], rtl)
    for i, (stat, label) in enumerate(t["opp_stats"]):
        x = Inches(0.65 + i * 3.15)
        add_rect(s, x, Inches(1.7), Inches(3.0), Inches(2.2), WHITE)
        add_textbox(s, x + Inches(0.15), Inches(1.9), Inches(2.7), Inches(0.7), stat, 28, True, BLUE, "Cambria", rtl=rtl, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(2.7), Inches(2.7), Inches(1.0), label, 12, False, DARK, rtl=rtl, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(4.3), Inches(11.8), Inches(1.4), t["opp_body"], 14, False, DARK, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} opportunity")

    # --- 7 Solution roles ---
    s = next_slide()
    header(s, t["sol_eye"], t["sol_title"], rtl)
    for i, (letter, role, desc) in enumerate(t["sol_roles"]):
        x = Inches(0.65 + i * 3.15)
        add_rect(s, x, Inches(1.7), Inches(3.0), Inches(3.8), WHITE)
        add_oval(s, x + Inches(1.0), Inches(2.0), Inches(1.0), Inches(1.0), BLUE)
        add_textbox(s, x + Inches(1.0), Inches(2.15), Inches(1.0), Inches(0.7), letter, 22, True, WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(3.2), Inches(2.7), Inches(0.4), role, 16, True, NAVY, rtl=rtl, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(3.7), Inches(2.7), Inches(1.5), desc, 12, False, DARK, rtl=rtl, align=PP_ALIGN.CENTER)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} solution")

    # --- 8 Features ---
    s = next_slide()
    header(s, t["feat_eye"], t["feat_title"], rtl)
    for i, (title, items) in enumerate(t["feat_cols"]):
        x = Inches(0.55 + i * 3.2)
        add_rect(s, x, Inches(1.65), Inches(3.05), Inches(4.0), WHITE)
        add_rect(s, x, Inches(1.65), Inches(3.05), Inches(0.55), BLUE)
        add_textbox(s, x + Inches(0.1), Inches(1.72), Inches(2.85), Inches(0.4), title, 13, True, WHITE, rtl=rtl, align=PP_ALIGN.CENTER)
        y = Inches(2.45)
        for it in items:
            add_textbox(s, x + Inches(0.15), y, Inches(2.75), Inches(0.45), "• " + it, 12, False, DARK, rtl=rtl, align=al)
            y += Inches(0.55)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} features")

    # --- 9 Development status ---
    s = next_slide()
    header(s, t["dev_eye"], t["dev_title"], rtl)
    add_rect(s, Inches(0.65), Inches(1.6), Inches(5.7), Inches(4.1), WHITE)
    add_rect(s, Inches(6.7), Inches(1.6), Inches(5.7), Inches(4.1), WHITE)
    add_textbox(s, Inches(0.85), Inches(1.75), Inches(5.3), Inches(0.4), "✓" if not rtl else "✓", 14, True, BLUE, rtl=rtl, align=al)
    done_label = "Réalisé / en place" if not rtl else "منجز / قائم"
    todo_label = "Restant à financer" if not rtl else "متبقي للتمويل"
    add_textbox(s, Inches(1.2), Inches(1.75), Inches(4.8), Inches(0.4), done_label, 14, True, BLUE, rtl=rtl, align=al)
    add_textbox(s, Inches(6.9), Inches(1.75), Inches(5.3), Inches(0.4), todo_label, 14, True, AMBER, rtl=rtl, align=al)
    y = Inches(2.3)
    for item in t["dev_done"]:
        add_textbox(s, Inches(0.9), y, Inches(5.2), Inches(0.35), "• " + item, 11, False, DARK, rtl=rtl, align=al)
        y += Inches(0.38)
    y = Inches(2.3)
    for item in t["dev_todo"]:
        add_textbox(s, Inches(6.9), y, Inches(5.2), Inches(0.35), "• " + item, 11, False, DARK, rtl=rtl, align=al)
        y += Inches(0.38)
    note_bar(s, t["dev_note"], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} status")

    # --- 10 Feasibility ---
    s = next_slide()
    header(s, t["fea_eye"], t["fea_title"], rtl)
    positions = [
        (Inches(0.65), Inches(1.7)), (Inches(6.7), Inches(1.7)),
        (Inches(0.65), Inches(3.9)), (Inches(6.7), Inches(3.9)),
    ]
    for (title, body), (x, y) in zip(t["fea_cards"], positions):
        card(s, x, y, Inches(5.7), Inches(1.9), title, body, BLUE, rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} feasibility")

    # --- 11 Technical ---
    s = next_slide()
    header(s, t["tech_eye"], t["tech_title"], rtl)
    for i, (num, title, body) in enumerate(t["tech_items"]):
        col, row = i % 3, i // 3
        x = Inches(0.65 + col * 4.1)
        y = Inches(1.7 + row * 2.1)
        add_rect(s, x, y, Inches(3.9), Inches(1.9), WHITE)
        add_oval(s, x + Inches(0.2), y + Inches(0.25), Inches(0.55), Inches(0.55), BLUE)
        add_textbox(s, x + Inches(0.2), y + Inches(0.32), Inches(0.55), Inches(0.45), num, 14, True, WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.9), y + Inches(0.3), Inches(2.8), Inches(0.4), title, 13, True, NAVY, rtl=rtl, align=al)
        add_textbox(s, x + Inches(0.2), y + Inches(0.95), Inches(3.5), Inches(0.8), body, 11, False, DARK, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} technical")

    # --- 12 Security ---
    s = next_slide()
    header(s, t["sec_eye"], t["sec_title"], rtl)
    for i, item in enumerate(t["sec_items"]):
        col, row = i % 2, i // 2
        x = Inches(0.65 + col * 6.2)
        y = Inches(1.7 + row * 1.3)
        add_rect(s, x, y, Inches(5.9), Inches(1.15), WHITE)
        add_rect(s, x, y, Inches(0.12), Inches(1.15), BLUE)
        add_textbox(s, x + Inches(0.35), y + Inches(0.35), Inches(5.3), Inches(0.5), item, 14, False, DARK, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} security")

    # --- 13 Content plan ---
    s = next_slide()
    header(s, t["cont_eye"], t["cont_title"], rtl)
    for i, (num, label) in enumerate(t["cont_steps"]):
        x = Inches(0.5 + i * 2.1)
        add_oval(s, x + Inches(0.55), Inches(1.75), Inches(0.7), Inches(0.7), BLUE)
        add_textbox(s, x + Inches(0.55), Inches(1.88), Inches(0.7), Inches(0.5), num, 14, True, WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(2.6), Inches(2.0), Inches(0.8), label, 11, True, NAVY, rtl=rtl, align=PP_ALIGN.CENTER)
        if i < 5:
            add_rect(s, x + Inches(1.7), Inches(2.0), Inches(0.5), Inches(0.06), PALE)
    for i, (stat, label) in enumerate(t["cont_stats"]):
        x = Inches(0.65 + i * 3.15)
        add_rect(s, x, Inches(3.7), Inches(3.0), Inches(1.9), NOTE_BG)
        add_textbox(s, x + Inches(0.1), Inches(3.9), Inches(2.8), Inches(0.6), stat, 24, True, BLUE, "Cambria", rtl=rtl, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(4.6), Inches(2.8), Inches(0.8), label, 12, False, DARK, rtl=rtl, align=PP_ALIGN.CENTER)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} content")

    # --- 14 Budget 1/2 ---
    s = next_slide()
    header(s, t["bud_eye"], t["bud1_title"], rtl)
    rows = [t["bud_headers"]]
    for (num, cost), label in zip(BUDGET_ROWS[:9], t["bud_labels"][:9]):
        rows.append([str(num), label, fmt_tnd(cost)])
    table_slide(s, rows, Inches(0.7), Inches(1.55), Inches(11.8), [Inches(0.8), Inches(8.5), Inches(2.5)], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} budget1")

    # --- 15 Budget 2/2 ---
    s = next_slide()
    header(s, t["bud_eye"], t["bud2_title"], rtl)
    rows = [t["bud_headers"]]
    for (num, cost), label in zip(BUDGET_ROWS[9:], t["bud_labels"][9:]):
        rows.append([str(num), label, fmt_tnd(cost)])
    total_label = "TOTAL PROJET" if not rtl else "إجمالي المشروع"
    rows.append(["", total_label, fmt_tnd(127000)])
    table_slide(s, rows, Inches(0.7), Inches(1.55), Inches(11.8), [Inches(0.8), Inches(8.5), Inches(2.5)], rtl)
    note = "Coût total du projet : 127 000 TND — dont 100 000 TND demandés au Ministère." if not rtl else \
           "إجمالي تكلفة المشروع: 127 000 د.ت — منها 100 000 د.ت مطلوبة من الوزارة."
    note_bar(s, note, rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} budget2")

    # --- 16 Funding ---
    s = next_slide()
    header(s, t["fund_eye"], t["fund_title"], rtl)
    table_slide(s, [t["fund_headers"]] + t["fund_rows"], Inches(1.5), Inches(2.0), Inches(10.3),
                [Inches(6.5), Inches(2.0), Inches(1.8)], rtl)
    # big numbers
    add_rect(s, Inches(0.65), Inches(4.5), Inches(5.7), Inches(1.2), NOTE_BG)
    add_rect(s, Inches(6.7), Inches(4.5), Inches(5.7), Inches(1.2), NOTE_BG)
    add_textbox(s, Inches(0.85), Inches(4.7), Inches(5.3), Inches(0.8),
                "100 000 TND  ·  79 %" if not rtl else "100 000 د.ت  ·  79 %", 22, True, BLUE, "Cambria", rtl=rtl, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.9), Inches(4.7), Inches(5.3), Inches(0.8),
                "27 000 TND  ·  21 %" if not rtl else "27 000 د.ت  ·  21 %", 22, True, DEEP, "Cambria", rtl=rtl, align=PP_ALIGN.CENTER)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} funding")

    # --- 17 Owner contribution ---
    s = next_slide()
    header(s, t["own_eye"], t["own_title"], rtl)
    table_slide(s, [t["own_headers"]] + t["own_rows"], Inches(2.0), Inches(1.8), Inches(9.3),
                [Inches(6.5), Inches(2.8)], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} owner")

    # --- 18 HR ---
    s = next_slide()
    header(s, t["hr_eye"], t["hr_title"], rtl)
    table_slide(s, [t["hr_headers"]] + t["hr_rows"], Inches(2.5), Inches(1.6), Inches(8.3),
                [Inches(6.0), Inches(2.3)], rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} hr")

    # --- 19 Timeline ---
    s = next_slide()
    header(s, t["tl_eye"], t["tl_title"], rtl)
    for i, (period, title, body) in enumerate(t["tl_phases"]):
        y = Inches(1.55 + i * 0.9)
        add_rect(s, Inches(0.65), y, Inches(11.9), Inches(0.8), WHITE)
        add_rect(s, Inches(0.65), y, Inches(0.12), Inches(0.8), BLUE if i % 2 == 0 else AMBER)
        add_textbox(s, Inches(0.95), y + Inches(0.1), Inches(2.2), Inches(0.6), period, 12, True, BLUE, rtl=rtl, align=al)
        add_textbox(s, Inches(3.2), y + Inches(0.1), Inches(2.5), Inches(0.6), title, 13, True, NAVY, rtl=rtl, align=al)
        add_textbox(s, Inches(5.8), y + Inches(0.1), Inches(6.5), Inches(0.6), body, 12, False, DARK, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} timeline")

    # --- 20 Risks ---
    s = next_slide()
    header(s, t["risk_eye"], t["risk_title"], rtl)
    for i, (kind, risk, mit) in enumerate(t["risk_items"]):
        col, row = i % 2, i // 2
        x = Inches(0.65 + col * 6.2)
        y = Inches(1.65 + row * 2.15)
        add_rect(s, x, y, Inches(5.95), Inches(2.0), WHITE)
        add_textbox(s, x + Inches(0.25), y + Inches(0.2), Inches(5.4), Inches(0.35), kind, 12, True, BLUE, rtl=rtl, align=al)
        add_textbox(s, x + Inches(0.25), y + Inches(0.6), Inches(5.4), Inches(0.4), risk, 14, True, NAVY, rtl=rtl, align=al)
        add_textbox(s, x + Inches(0.25), y + Inches(1.1), Inches(5.4), Inches(0.7), mit, 12, False, DARK, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} risks")

    # --- 21 Sustainability ---
    s = next_slide()
    header(s, t["sus_eye"], t["sus_title"], rtl)
    positions = [
        (Inches(0.65), Inches(1.7)), (Inches(6.7), Inches(1.7)),
        (Inches(0.65), Inches(3.9)), (Inches(6.7), Inches(3.9)),
    ]
    for (title, body), (x, y) in zip(t["sus_cards"], positions):
        card(s, x, y, Inches(5.7), Inches(1.9), title, body, BLUE, rtl)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} sustainability")

    # --- 22 KPIs ---
    s = next_slide()
    header(s, t["kpi_eye"], t["kpi_title"], rtl)
    for i, (stat, label) in enumerate(t["kpi_items"]):
        col, row = i % 3, i // 3
        x = Inches(0.65 + col * 4.1)
        y = Inches(1.7 + row * 2.2)
        add_rect(s, x, y, Inches(3.9), Inches(2.0), WHITE)
        add_textbox(s, x + Inches(0.15), y + Inches(0.35), Inches(3.6), Inches(0.7), stat, 26, True, BLUE, "Cambria", rtl=rtl, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), y + Inches(1.15), Inches(3.6), Inches(0.7), label, 12, False, DARK, rtl=rtl, align=PP_ALIGN.CENTER)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} kpis")

    # --- 23-26 Annex screenshots ---
    annexes = [
        (t["ax_home"], ["home1_0.png", "home1_1.png", "home2_0.png"]),
        (t["ax_admin"], ["admin_0.png"]),
        (t["ax_student"], ["student_0.png"]),
        (t["ax_teacher"], ["teacher_0.png"]),
    ]
    for title, files in annexes:
        s = next_slide()
        header(s, t["ax_eye"], title, rtl)
        imgs = [MEDIA / f for f in files if (MEDIA / f).exists()]
        if not imgs:
            # fallback older names
            imgs = list(MEDIA.glob("slide19*.png"))[:2]
        if len(imgs) == 1:
            s.shapes.add_picture(str(imgs[0]), Inches(1.5), Inches(1.6), height=Inches(4.8))
        elif len(imgs) >= 2:
            s.shapes.add_picture(str(imgs[0]), Inches(0.5), Inches(1.6), width=Inches(6.0))
            s.shapes.add_picture(str(imgs[1]), Inches(6.7), Inches(1.6), width=Inches(6.0))
        footer(s, n, t["brand"], rtl)
        log(f"  slide {n} annex {title[:20]}")

    # --- 27 Logo annex ---
    s = next_slide()
    header(s, t["ax_eye"], t["ax_logo"], rtl)
    if (MEDIA / "logo.jpg").exists():
        s.shapes.add_picture(str(MEDIA / "logo.jpg"), Inches(5.15), Inches(2.2), Inches(3.0), Inches(3.0))
    add_textbox(s, Inches(2), Inches(5.4), Inches(9.3), Inches(0.4), t["brand_ar_line"], 16, True, NAVY, rtl=True, align=PP_ALIGN.CENTER)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} logo")

    # --- 28 Closing ---
    s = next_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_oval(s, Inches(9.6), Inches(-2.2), Inches(6.5), Inches(6.5), BLUE)
    add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5), t["brand"], 14, True, PALE, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(2.2), Inches(11), Inches(1.0), t["close_title"], 28, True, WHITE, "Cambria", rtl=rtl, align=al)
    add_rect(s, Inches(0.8), Inches(3.5), Inches(6.5), Inches(0.55), WHITE)
    add_textbox(s, Inches(0.95), Inches(3.55), Inches(6.2), Inches(0.45), t["close_budget"], 14, True, BLUE, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(4.25), Inches(10), Inches(0.4), t["close_total"], 14, False, PALE, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(4.75), Inches(11), Inches(0.5), t["close_status"], 13, True, AMBER, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(5.5), Inches(10), Inches(0.3), t["close_portal"], 12, False, MUTED, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(5.9), Inches(10), Inches(0.3), t["close_contact"], 12, False, MUTED, rtl=rtl, align=al)
    add_textbox(s, Inches(0.8), Inches(6.4), Inches(10), Inches(0.3), t["close_min"], 11, False, MUTED, rtl=rtl, align=al)
    footer(s, n, t["brand"], rtl)
    log(f"  slide {n} closing — total slides={n}")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    log(f"BUILD done lang={lang} slides={n} file={out} size={out.stat().st_size}")
    return n


def main():
    log("=== MAIN start ===")
    # ensure media
    if not (MEDIA / "logo.jpg").exists():
        log("ERROR missing logo.jpg in /tmp/pptx_media")
        sys.exit(1)
    n_fr = build("fr", OUT_FR)
    n_ar = build("ar", OUT_AR)
    # marker for completion
    Path("/tmp/funding_pptx/DONE").write_text(f"fr={n_fr} ar={n_ar}\n", encoding="utf-8")
    log(f"=== MAIN complete FR={n_fr} AR={n_ar} ===")
    print(f"OK FR={OUT_FR} ({n_fr} slides)")
    print(f"OK AR={OUT_AR} ({n_ar} slides)")


if __name__ == "__main__":
    main()
